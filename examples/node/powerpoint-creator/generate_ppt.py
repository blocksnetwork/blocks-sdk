from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
import argparse
import json
from pathlib import Path

# --- PubNub Color Palette ---
PN_LIGHT_BG = RGBColor(0xF9, 0xF9, 0xF9) # Light Background
PN_DARK_BG = RGBColor(0x07, 0x0F, 0x39)  # Navy Blue
PN_RED = RGBColor(0xC7, 0x19, 0x29)      # PubNub Red
PN_BLUE = RGBColor(0x52, 0x8D, 0xFA)     # Light Blue
PN_TEXT = RGBColor(0x17, 0x17, 0x17)     # Very Dark Grey
PN_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Theme palette for JSON-driven presentations
THEMES = {
    "light": {
        "background": PN_LIGHT_BG,
        "title_bg": PN_DARK_BG,
        "title_text": PN_WHITE,
        "text": PN_TEXT,
        "accent": PN_RED,
        "secondary": PN_BLUE,
    },
    "dark": {
        "background": PN_DARK_BG,
        "title_bg": PN_DARK_BG,
        "title_text": PN_WHITE,
        "text": PN_WHITE,
        "accent": PN_RED,
        "secondary": PN_BLUE,
    },
    "corporate": {
        "background": PN_WHITE,
        "title_bg": PN_BLUE,
        "title_text": PN_WHITE,
        "text": PN_TEXT,
        "accent": PN_BLUE,
        "secondary": PN_RED,
    },
}

def get_theme(name: str):
    return THEMES.get(name, THEMES["light"])

def add_block(slide, x, y, w, h, text, bg_color, text_color=PN_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    line = shape.line
    line.color.rgb = PN_BLUE
    line.width = Pt(1.5)

    set_text_style(shape, text, font_size=Pt(14), color=text_color, bold=False, alignment=PP_ALIGN.CENTER)

    return shape

def add_connector(slide, shape1, shape2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, 0, 0)
    connector.begin_connect(shape1, 2) # Bottom
    connector.end_connect(shape2, 0)   # Top
    line = connector.line
    line.width = Pt(2)
    line.color.rgb = PN_TEXT
    return connector

def set_text_style(shape, text, font_size=Pt(18), color=PN_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    """Explicitly sets text with styling on a shape's text frame."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True # Ensure wrapping is ON

    p = tf.paragraphs[0]
    p.alignment = alignment

    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = 'Arial'

def apply_slide_style(slide, title_text, is_title_slide=False):
    """Applies common styling to slides."""

    # Set Background
    background = slide.background
    fill = background.fill
    fill.solid()

    if is_title_slide:
        fill.fore_color.rgb = PN_DARK_BG
        # Add bottom strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), Inches(13.333), Inches(0.3))
        strip.fill.solid()
        strip.fill.fore_color.rgb = PN_RED
        strip.line.fill.background() # No line
    else:
        fill.fore_color.rgb = PN_LIGHT_BG

        # Style Title using a Text Box explicitly to avoid placeholder issues
        # Remove default title placeholder if exists to avoid conflict
        if slide.shapes.title:
            sp = slide.shapes.title
            sp.element.getparent().remove(sp.element)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        set_text_style(title_box, title_text, font_size=Pt(36), color=PN_DARK_BG, bold=True)

        # Add Header/Footer accents
        # Red bar top left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(1.5), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PN_RED
        bar.line.fill.background()

        # Logo/Brand text bottom right
        textbox = slide.shapes.add_textbox(Inches(11), Inches(7), Inches(2), Inches(0.4))
        set_text_style(textbox, "PubNub Blocks Network", font_size=Pt(10), color=PN_DARK_BG, alignment=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# JSON-driven presentation helpers
# ---------------------------------------------------------------------------

def apply_slide_style_custom(slide, title_text, theme, is_title_slide=False):
    """Applies themed styling to slides for JSON-driven decks."""
    background = slide.background
    fill = background.fill
    fill.solid()

    if is_title_slide:
        fill.fore_color.rgb = theme["title_bg"]
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), Inches(13.333), Inches(0.3))
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme["accent"]
        strip.line.fill.background()
        return

    fill.fore_color.rgb = theme["background"]

    if slide.shapes.title:
        sp = slide.shapes.title
        sp.element.getparent().remove(sp.element)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    set_text_style(title_box, title_text, font_size=Pt(34), color=theme["text"], bold=True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(1.5), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

def add_bullets(text_frame, bullets, color, font_size=Pt(20)):
    text_frame.clear()
    text_frame.word_wrap = True
    for i, b in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = str(b)
        p.level = 1
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_before = Pt(6)

def add_body_text(text_frame, content, color, font_size=Pt(20)):
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = str(content)
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = 'Arial'

def create_presentation_from_spec(spec, output_path):
    """Create a PPTX from a JSON spec."""
    if not isinstance(spec, dict):
        raise ValueError("Presentation spec must be a JSON object.")

    title = str(spec.get("title", "Untitled Presentation"))
    subtitle = spec.get("subtitle")
    author = spec.get("author")
    theme = get_theme(str(spec.get("theme", "light")))
    slides_spec = spec.get("slides")
    if not isinstance(slides_spec, list):
        slides_spec = []

    prs = Presentation()

    # Set 16:9 Aspect Ratio explicitly
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_style_custom(slide, "", theme, is_title_slide=True)

    tb_title = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11), Inches(1.6))
    set_text_style(tb_title, title, font_size=Pt(52), color=theme["title_text"], bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        tb_sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.8))
        set_text_style(tb_sub, str(subtitle), font_size=Pt(24), color=theme["secondary"], alignment=PP_ALIGN.CENTER)

    if author:
        tb_author = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
        set_text_style(tb_author, f"Author: {author}", font_size=Pt(14), color=theme["secondary"], alignment=PP_ALIGN.CENTER)

    # --- Content Slides ---
    for slide_spec in slides_spec:
        if not isinstance(slide_spec, dict):
            continue
        slide_title = str(slide_spec.get("title", ""))
        layout = slide_spec.get("layout")
        bullets = slide_spec.get("bullets")
        content = slide_spec.get("content", "")
        notes = slide_spec.get("notes")

        if layout is None:
            layout = "bullets" if isinstance(bullets, list) else "content"

        slide = prs.slides.add_slide(slide_layout)
        apply_slide_style_custom(slide, slide_title, theme, is_title_slide=False)

        if layout == "two-column" and isinstance(bullets, list):
            half = (len(bullets) + 1) // 2
            left = bullets[:half]
            right = bullets[half:]

            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(5))
            add_bullets(left_box.text_frame, left, theme["text"])

            right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(6), Inches(5))
            add_bullets(right_box.text_frame, right, theme["text"])
        elif layout == "bullets" and isinstance(bullets, list):
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(5))
            add_bullets(content_box.text_frame, bullets, theme["text"])
        else:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(5))
            add_body_text(content_box.text_frame, content, theme["text"])

        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    outpath = Path(output_path)
    outpath.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(outpath))
    print(f"Presentation created successfully at {outpath}")


def main():
    parser = argparse.ArgumentParser(description="Generate a PowerPoint deck.")
    parser.add_argument(
        "--input",
        help="Path to JSON presentation spec.",
    )
    parser.add_argument(
        "--output",
        default="output.pptx",
        help="Output PPTX path (default: output.pptx).",
    )
    args = parser.parse_args()

    if args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            spec = json.load(f)
        create_presentation_from_spec(spec, args.output)
    else:
        print("Error: --input is required. Provide a JSON presentation spec.")
        parser.print_help()

if __name__ == "__main__":
    main()
